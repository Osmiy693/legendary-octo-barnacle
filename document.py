"""
document.py - Document class hierarchy for AI Study Assistant.

Demonstrates OOP principles:
- ENCAPSULATION: _title and _content are protected attributes, accessed via getters.
- INHERITANCE: TextDocument, PDFDocument, DocxDocument, XlsxDocument, PptxDocument,
  and MultiFileDocument inherit from Document.
- POLYMORPHISM: process() is overridden in each subclass.
"""

from pypdf import PdfReader
from docx import Document as DocxReader
from openpyxl import load_workbook
from pptx import Presentation
import io
import os
from collections import defaultdict


class Document:
    """
    Parent class representing any study material document.
    Uses encapsulation to protect _title and _content attributes.
    Demonstrates: ENCAPSULATION, POLYMORPHISM (process() to be overridden).
    """

    def __init__(self, title: str = "Untitled"):
        self._title = title
        self._content = ""

    def get_title(self) -> str:
        """Getter for the protected _title attribute."""
        return self._title

    def get_content(self) -> str:
        """Getter for the protected _content attribute."""
        return self._content

    def process(self, data) -> None:
        """
        POLYMORPHISM: Base implementation.
        Subclasses override this to provide their own processing logic.
        """
        raise NotImplementedError("Subclasses must implement process() method")

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}(title='{self._title}', chars={len(self._content)})"


class TextDocument(Document):
    """
    INHERITANCE: Inherits from Document.
    Processes plain text input (from text area or .txt file).
    Demonstrates: INHERITANCE, POLYMORPHISM (overrides process()).
    """

    def process(self, data: str) -> None:
        """POLYMORPHISM: Stores the raw text directly as content."""
        if not data or not data.strip():
            raise ValueError("Cannot process empty text. Please provide some content.")
        self._content = data.strip()


class PDFDocument(Document):
    """
    INHERITANCE: Inherits from Document.
    Extracts text from PDF files using pypdf.
    Demonstrates: INHERITANCE, POLYMORPHISM (overrides process()).
    """

    def process(self, data: bytes) -> None:
        """POLYMORPHISM: Extracts text from PDF bytes using pypdf library."""
        if not data:
            raise ValueError("Cannot process empty PDF data.")
        try:
            reader = PdfReader(io.BytesIO(data))
            extracted_text = []
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    extracted_text.append(page_text)
            text = "\n\n".join(extracted_text)
            if not text.strip():
                raise ValueError(
                    "No extractable text found in this PDF. "
                    "The file may be scanned or image-based."
                )
            self._content = text.strip()
        except Exception as e:
            raise ValueError(f"Failed to process PDF file: {str(e)}")


class DocxDocument(Document):
    """
    INHERITANCE: Inherits from Document.
    Extracts text from .docx files using python-docx.
    """

    def process(self, data: bytes) -> None:
        """POLYMORPHISM: Extracts text from DOCX bytes."""
        if not data:
            raise ValueError("Cannot process empty DOCX data.")
        try:
            doc = DocxReader(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            if not text.strip():
                # Also try tables
                tables_text = []
                for table in doc.tables:
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            tables_text.append(" | ".join(cells))
                text = "\n".join(tables_text)
            if not text.strip():
                raise ValueError("No extractable text found in this DOCX.")
            self._content = text.strip()
        except Exception as e:
            raise ValueError(f"Failed to process DOCX file: {str(e)}")


class XlsxDocument(Document):
    """
    INHERITANCE: Inherits from Document.
    Extracts text from .xlsx/.xls files using openpyxl.
    """

    def process(self, data: bytes) -> None:
        """POLYMORPHISM: Extracts text from Excel bytes."""
        if not data:
            raise ValueError("Cannot process empty Excel data.")
        try:
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            all_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_lines = [f"--- Sheet: {sheet_name} ---"]
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    line = " | ".join(c for c in cells if c)
                    if line.strip():
                        sheet_lines.append(line)
                all_text.extend(sheet_lines)
            wb.close()
            text = "\n".join(all_text)
            if not text.strip():
                raise ValueError("No extractable text found in this Excel file.")
            self._content = text.strip()
        except Exception as e:
            raise ValueError(f"Failed to process Excel file: {str(e)}")


class PptxDocument(Document):
    """
    INHERITANCE: Inherits from Document.
    Extracts text from .pptx files using python-pptx.
    """

    def process(self, data: bytes) -> None:
        """POLYMORPHISM: Extracts text from PPTX bytes."""
        if not data:
            raise ValueError("Cannot process empty PPTX data.")
        try:
            prs = Presentation(io.BytesIO(data))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_lines = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_lines.append(para.text.strip())
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_lines.append(" | ".join(cells))
                slides_text.append("\n".join(slide_lines))
            text = "\n\n".join(slides_text)
            if not text.strip():
                raise ValueError("No extractable text found in this PPTX.")
            self._content = text.strip()
        except Exception as e:
            raise ValueError(f"Failed to process PPTX file: {str(e)}")


class MultiFileDocument(Document):
    """
    INHERITANCE: Inherits from Document.
    Holds multiple files (from folder upload) and merges their content.
    Also maintains a file tree for Finder-like display.
    Demonstrates: INHERITANCE, POLYMORPHISM (overrides process()), COMPOSITION.
    """

    def __init__(self, title: str = "Uploaded Folder"):
        super().__init__(title)
        self._files = []
        self._file_tree = {}

    def add_file(self, filename: str, content: str) -> None:
        """Add a single file's content to the combined document."""
        self._files.append({
            "name": filename,
            "content": content,
            "chars": len(content)
        })

    def get_files(self) -> list:
        """Returns list of all files with their metadata."""
        return self._files

    def get_file_tree(self) -> dict:
        """Build and return a nested directory tree from file paths."""
        tree = {}
        for f in self._files:
            parts = f["name"].replace("\\", "/").split("/")
            current = tree
            for i, part in enumerate(parts):
                if i == len(parts) - 1:
                    if "__files__" not in current:
                        current["__files__"] = []
                    current["__files__"].append({
                        "name": part,
                        "chars": f["chars"],
                        "full_path": f["name"]
                    })
                else:
                    if part not in current:
                        current[part] = {}
                    current = current[part]
        return tree

    def get_file_count(self) -> int:
        """Returns the total number of files."""
        return len(self._files)

    def get_total_chars(self) -> int:
        """Returns combined character count of all files."""
        return sum(f["chars"] for f in self._files)

    def process(self, data=None) -> None:
        """
        POLYMORPHISM: Combines all file contents with headers.
        """
        if not self._files:
            raise ValueError("No files have been added to the folder document.")
        combined = []
        for f in self._files:
            combined.append(f"--- FILE: {f['name']} ---\n{f['content']}")
        self._content = "\n\n".join(combined)
